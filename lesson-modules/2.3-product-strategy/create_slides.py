#!/usr/bin/env python3
"""Create executive presentation from TaskFlow AI Strategy document."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """Add title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(100, 100, 100)

    return slide

def add_section_slide(prs, section_title):
    """Add section header slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8), Inches(10), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(47, 84, 150)
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, subbullets=None):
    """Add content slide with bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(47, 84, 150)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)

        # Add subbullets if provided
        if subbullets and i < len(subbullets) and subbullets[i]:
            for sub in subbullets[i]:
                p = tf.add_paragraph()
                p.text = f"    ◦ {sub}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(80, 80, 80)
                p.space_after = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(47, 84, 150)

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.space_after = Pt(10)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.space_after = Pt(10)

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add slide with table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(47, 84, 150)

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * num_rows)).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(47, 84, 150)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)

    return slide

def create_presentation():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
                    "TaskFlow H1 2026\nAI Product Strategy",
                    "Gen AI PM | Декабрь 2024")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Конкуренты уходят в enterprise — SMB недообслужен",
        "Наш подход: партнёрство для AI, focused tools, move fast",
        "AI включён в базовую цену $12/мес — ставка на retention",
        "H1 2026: 3 AI-инструмента (Meeting Notes, Task Breakdown, Standup)",
        "Target: 50% WAU using AI, +10% LTV"
    ])

    # Slide 3: Section - Diagnosis
    add_section_slide(prs, "ДИАГНОЗ")

    # Slide 4: Competitive Landscape
    add_table_slide(prs, "Конкурентный ландшафт",
        ["Конкурент", "AI-подход", "Ценообразование", "Фокус"],
        [
            ["Notion", "AI Assistant + Agents", "От $20/мес", "Enterprise"],
            ["Linear", "AI для инженеров", "Включено", "Dev teams"],
            ["Asana", "AI Studio", "От $25/мес", "Enterprise"],
            ["TaskFlow", "Focused AI tools", "$12 всё вкл.", "SMB"]
        ])

    # Slide 5: Our Position
    add_two_column_slide(prs, "Позиция TaskFlow",
        "Сильные стороны",
        ["Voice chat работает — первопроходцы",
         "Activation 56% для SMB",
         "Понимаем маленькие команды"],
        "Ограничения",
        ["2 AI-инженера",
         "$50k/квартал бюджет",
         "$3/user AI-косты"])

    # Slide 6: Core Challenge
    add_content_slide(prs, "Ядро стратегического вызова", [
        "Как TaskFlow может победить в AI с ограниченными ресурсами?",
        "Конкуренты имеют больше денег и инженеров",
        "Но все уходят в enterprise...",
        "SMB остаётся недообслуженным",
        "Наша возможность: focused AI для маленьких команд"
    ])

    # Slide 7: Section - Guiding Policy
    add_section_slide(prs, "РУКОВОДЯЩАЯ ПОЛИТИКА")

    # Slide 8: Strategic Direction
    add_content_slide(prs, "Наш стратегический подход", [
        "ПАРТНЁРСТВО: Используем OpenAI/Anthropic, не строим свои модели",
        "ФОКУС: Свой roadmap, не реагируем на конкурентов",
        "PRICING: AI в базовой цене $12 — ставка на retention",
        "SCOPE: Focused tools для конкретных jobs-to-be-done",
        "SPEED: Move fast, learn fast, iterate"
    ])

    # Slide 9: What We're NOT Doing
    add_content_slide(prs, "Чего мы НЕ делаем (tradeoffs)", [
        "Не строим свои AI-модели — не наш бизнес",
        "Не идём в enterprise — там Notion/Asana",
        "Не реагируем на каждый ход конкурентов",
        "Не делаем AI premium tier — SMB хотят простоту",
        "Не полируем месяцами — скорость важнее"
    ])

    # Slide 10: Section - Coherent Actions
    add_section_slide(prs, "СОГЛАСОВАННЫЕ ДЕЙСТВИЯ")

    # Slide 11: Q1 Roadmap
    add_content_slide(prs, "Q1 2026 Roadmap", [
        "ЯНВАРЬ: OpenAI интеграция + Feature flags",
        "ФЕВРАЛЬ: AI Meeting Notes MVP → 10% beta",
        "МАРТ: Meeting Notes → 100% rollout"
    ], subbullets=[
        ["GPT-4o для core AI", "Инфраструктура для экспериментов"],
        ["Voice → Summary → Tasks", "2-недельные циклы итераций"],
        ["Измерение retention impact", "Начало Task Breakdown"]
    ])

    # Slide 12: Q2 Roadmap
    add_content_slide(prs, "Q2 2026 Roadmap", [
        "АПРЕЛЬ: AI Task Breakdown MVP",
        "МАЙ: Task Breakdown → 50% + Daily Standup начало",
        "ИЮНЬ: Full rollout + H1 ретроспектива"
    ], subbullets=[
        ["Большая задача → подзадачи с estimates"],
        ["Cost optimization", "AI Standup Summary MVP"],
        ["3 AI tools shipped", "H2 планирование"]
    ])

    # Slide 13: Success Metrics
    add_table_slide(prs, "Success Metrics",
        ["Метрика", "Baseline", "Q1 Target", "Q2 Target"],
        [
            ["AI Feature Usage (WAU)", "0%", "30%", "50%"],
            ["Retention (Week 4)", "69%", "75%", "80%"],
            ["AI Error Rate", "—", "<5%", "<3%"],
            ["LTV Impact", "—", "—", "+10%"]
        ])

    # Slide 14: Why We'll Win
    add_content_slide(prs, "Почему мы победим", [
        "ПРОСТОТА: $12/мес всё включено vs Notion $20+",
        "ФОКУС: Конкретные tools vs general AI assistant",
        "СКОРОСТЬ: Маленькая команда = быстрые решения",
        "SMB-ЭКСПЕРТИЗА: Понимаем workflows маленьких команд",
        "ИНТЕГРАЦИЯ: AI внутри workflow, не отдельный чат-бот"
    ])

    # Slide 15: The Ask
    add_content_slide(prs, "The Ask", [
        "ОТ ENGINEERING: 2 инженера focused на AI весь H1",
        "ОТ LEADERSHIP: Approval на margin hit (~25% reduction)",
        "ОТ LEADERSHIP: Patience — ROI через retention, не immediate",
        "ОТ DESIGN: Shared resources на AI tools UX"
    ])

    # Save
    output_path = "/Users/lance/Documents/claude-code-course/lesson-modules/2.3-product-strategy/strategy-review-slides.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
